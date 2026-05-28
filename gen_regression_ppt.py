"""
Generate: GFX_GOP_Regression_Checker_Overview.pptx  (2-slide edition)
Run with the project venv: python gen_regression_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# __ Palette ___________________________________________________________________
INTEL_BLUE   = RGBColor(0x00, 0x68, 0xB5)
INTEL_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GOLD  = RGBColor(0xD4, 0xA0, 0x17)
ACCENT_GREEN = RGBColor(0x16, 0xA3, 0x4A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY     = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY    = RGBColor(0x37, 0x41, 0x51)
PALE_BLUE    = RGBColor(0xEF, 0xF6, 0xFF)
PALE_GOLD    = RGBColor(0xFF, 0xFB, 0xEB)
PALE_GREEN   = RGBColor(0xF0, 0xFD, 0xF4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# __ Helpers ___________________________________________________________________

def rect(slide, l, t, w, h, fill=None, border=None, bpt=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border and bpt:
        s.line.color.rgb = border; s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s


def label(slide, text, l, t, w, h, size=10, bold=False, color=DARK_GRAY,
          align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tx.word_wrap = True
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bool(bold)
    r.font.color.rgb = color
    return tx


def header(slide, title, sub, accent=INTEL_BLUE):
    rect(slide, 0, 0, 13.33, 0.9, fill=accent)
    label(slide, title, 0.3, 0.04, 10, 0.5, size=24, bold=True, color=WHITE)
    label(slide, sub, 0.3, 0.54, 12, 0.32, size=10,
          color=RGBColor(0xBF, 0xDB, 0xFF))


def footer(slide, txt="GFX / GOP Driver Regression Checker"):
    rect(slide, 0, 7.22, 13.33, 0.28, fill=LIGHT_GRAY)
    label(slide, txt, 0.3, 7.24, 12.7, 0.24, size=7.5, color=MID_GRAY)


# ==============================================================================
# SLIDE 1  --  Overall Workflow
# ==============================================================================

s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, 13.33, 7.5, fill=WHITE)
header(s1,
       "GFX / GOP Driver Regression Checker -- Workflow",
       "What happens when you click  Check GFX Regression  or  Check GOP Regression")
footer(s1)

# __ Step flow (8 boxes) _______________________________________________________
steps = [
    ("1\nRead HSD ID",             INTEL_BLUE),
    ("2\nPOST\n/regression-check", INTEL_BLUE),
    ("3\nHSD API\nFetch bug data",  INTEL_BLUE),
    ("4\nVersions\ncomplete?",      ACCENT_GOLD),
    ("5\n/qb-check-auth\nAuth?",    INTEL_DARK),
    ("6\nPOST\n/qb-builds",        INTEL_BLUE),
    ("7\nResolve CI\nbuild ref",    INTEL_BLUE),
    ("8\nRender\nResults",          ACCENT_GREEN),
]

bw, bh = 1.42, 0.88
gap = 0.15
for i, (txt, col) in enumerate(steps):
    x = 0.25 + i * (bw + gap)
    rect(s1, x, 1.0, bw, bh, fill=col)
    label(s1, txt, x, 1.0, bw, bh, size=8.5, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        label(s1, ">", x + bw, 1.25, gap + 0.02, 0.35,
              size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Branch callouts
rect(s1, 0.25, 2.0, 13.0, 0.38, fill=PALE_GOLD, border=ACCENT_GOLD, bpt=1)
label(s1,
      "Step 4 -- versions missing: Manual Version Input Overlay shown     "
      "Step 5 -- not logged in: QB Login form shown before search",
      0.4, 2.04, 12.6, 0.3, size=9, color=RGBColor(0x78, 0x35, 0x00))

# __ 4 Network round-trips _____________________________________________________
rect(s1, 0.25, 2.5, 13.0, 0.28, fill=INTEL_BLUE)
label(s1, "4 Key Network Round-Trips (normal path)",
      0.4, 2.52, 12.5, 0.24, size=10, bold=True, color=WHITE)

trips = [
    ("(1) HSD API",         "Fetch bug metadata\n(title, status, versions)",          INTEL_BLUE),
    ("(2) QB Auth Check",   "Quick /rest/version ping\nSkip login if cached",          INTEL_DARK),
    ("(3) QB Build Search", "Match fail/pass/fix versions\nResolve CI build ref",      INTEL_BLUE),
    ("(4) QB Commits",      "On-demand: trace revenue-pr\n-> commit list",             ACCENT_GREEN),
]
for i, (title, desc, col) in enumerate(trips):
    tx = 0.3 + i * 3.25
    rect(s1, tx, 2.83, 3.1, 0.28, fill=col)
    label(s1, title, tx, 2.83, 3.1, 0.28,
          size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s1, tx, 3.11, 3.1, 0.58, fill=PALE_BLUE, border=col, bpt=0.8)
    label(s1, desc, tx + 0.05, 3.13, 3.0, 0.54, size=9, color=DARK_GRAY)

# __ Optional commit flow ______________________________________________________
rect(s1, 0.25, 3.82, 13.0, 0.28, fill=ACCENT_GREEN)
label(s1,
      "Optional -- Show Commits Flow  (triggered by Show Commits button)",
      0.4, 3.84, 12.5, 0.24, size=10, bold=True, color=WHITE)

commit_steps = [
    "User clicks\nShow Commits",
    "POST\n/qb-commits",
    "Extract\nrevenue-pr-NNNNN",
    "Search QB for\nlist-changes build",
    "GET /rest/builds\n/{id}/changes",
    "Parse XML\ncommit list",
    "Render table\n+ QB link",
]
bw2, bh2 = 1.72, 0.7
for i, txt in enumerate(commit_steps):
    cx = 0.28 + i * (bw2 + 0.1)
    rect(s1, cx, 4.16, bw2, bh2, fill=PALE_GREEN, border=ACCENT_GREEN, bpt=0.8)
    label(s1, txt, cx, 4.16, bw2, bh2, size=8.5, color=DARK_GRAY,
          align=PP_ALIGN.CENTER)
    if i < len(commit_steps) - 1:
        label(s1, ">", cx + bw2, 4.35, 0.12, 0.3,
              size=12, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# __ UI components map _________________________________________________________
rect(s1, 0.25, 5.0, 13.0, 0.24, fill=INTEL_DARK)
label(s1, "UI & Backend Components",
      0.4, 5.02, 12.5, 0.2, size=9, bold=True, color=WHITE)

comps = [
    ("sidepanel.html",      "Layout & buttons"),
    ("sidepanel.js",        "Mode switch\nkeydown (C=clear)"),
    ("regression_ui.js",    "checkGfxRegression()\ncheckGopRegression()\nshowQbLoginOverlay()"),
    ("bridge_server.py",    "HTTP routing\nport 8776"),
    ("regression_checker.py", "All HSD + QB\nbusiness logic"),
]
cw = 2.55
for i, (name, desc) in enumerate(comps):
    cx = 0.28 + i * (cw + 0.06)
    rect(s1, cx, 5.28, cw, 0.24, fill=INTEL_BLUE)
    label(s1, name, cx, 5.28, cw, 0.24,
          size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s1, cx, 5.52, cw, 0.7, fill=PALE_BLUE, border=INTEL_BLUE, bpt=0.6)
    label(s1, desc, cx + 0.04, 5.54, cw - 0.08, 0.66, size=8, color=DARK_GRAY)


# ==============================================================================
# SLIDE 2  --  Mermaid-Style Workflow Diagram (English)
# ==============================================================================

# -- Mermaid-style node helpers -------------------------------------------------
def pnode(sl, text, l, t, w, h, fill=None, border=None, bpt=1.2,
          size=9, bold=False, fc=None):
    """Rounded-rectangle process node (autoshape 5)."""
    _f = fill  or RGBColor(0xE0, 0xE7, 0xFF)
    _b = border or RGBColor(0x4F, 0x46, 0xE5)
    _c = fc    or RGBColor(0x1E, 0x29, 0x3B)
    s = sl.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = _f
    s.line.color.rgb = _b; s.line.width = Pt(bpt)
    tx = sl.shapes.add_textbox(Inches(l+0.06), Inches(t+0.04),
                               Inches(w-0.10), Inches(h-0.07))
    tx.word_wrap = True; tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bool(bold); r.font.color.rgb = _c

def dnode(sl, text, l, t, w, h, fill=None, border=None, bpt=1.2,
          size=8.5, fc=None):
    """Diamond decision node (autoshape 4)."""
    _f = fill  or RGBColor(0xFE, 0xF3, 0xC7)
    _b = border or RGBColor(0xD9, 0x77, 0x06)
    _c = fc    or RGBColor(0x1E, 0x29, 0x3B)
    s = sl.shapes.add_shape(4, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = _f
    s.line.color.rgb = _b; s.line.width = Pt(bpt)
    # Textbox in the centre third (diamond safe zone)
    tx = sl.shapes.add_textbox(Inches(l + w*0.22), Inches(t+0.06),
                               Inches(w*0.56),    Inches(h-0.12))
    tx.word_wrap = True; tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = _c

s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, 13.33, 7.5, fill=WHITE)
header(s2, "Overall Workflow",
       "Check GFX / GOP Driver Regression  �  End-to-End Execution Flow",
       accent=RGBColor(0x1E, 0x29, 0x3B))
footer(s2, "GFX / GOP Driver Regression Checker � Workflow")

# -- Mermaid color aliases ------------------------------------------------------
_PROC_F  = RGBColor(0xE0, 0xE7, 0xFF)   # light indigo  � process
_PROC_B  = RGBColor(0x4F, 0x46, 0xE5)   # indigo        � process border
_DEC_F   = RGBColor(0xFE, 0xF3, 0xC7)   # light amber   � decision
_DEC_B   = RGBColor(0xD9, 0x77, 0x06)   # amber
_START_F = RGBColor(0x4F, 0x46, 0xE5)   # filled indigo � trigger
_END_F   = RGBColor(0x16, 0xA3, 0x4A)   # filled green  � end result
_BYP_F   = RGBColor(0xFF, 0xF7, 0xED)   # light orange  � bypass branch
_BYP_B   = RGBColor(0xEA, 0x58, 0x0C)   # orange
_AUTH_F  = RGBColor(0xEF, 0xF6, 0xFF)   # light blue    � auth branch
_AUTH_B  = RGBColor(0x25, 0x63, 0xEB)   # blue
_CMT_F   = RGBColor(0xF0, 0xFD, 0xF4)   # light green   � commits
_CMT_B   = RGBColor(0x16, 0xA3, 0x4A)   # green
_LINE    = RGBColor(0x94, 0xA3, 0xB8)   # slate         � connectors
_TXT_D   = RGBColor(0x1E, 0x29, 0x3B)   # near-black

# -- Layout parameters ----------------------------------------------------------
MX, MW = 3.1, 5.0    # main column x, width
BX, BW = 0.1, 2.8    # left bypass branch
RX     = 8.3         # right branch x
BH     = 0.44        # box height (all nodes)
SP     = 0.54        # step spacing (top-to-top)
Y2     = [0.97 + i * SP for i in range(11)]

# -- Main flow nodes ------------------------------------------------------------
# 0: Trigger (filled indigo, white text)
pnode(s2, "User clicks   Check GFX / GOP Regression",
      MX, Y2[0], MW, BH, fill=_START_F, border=_START_F, bold=True, fc=WHITE)

# 1-3: Process steps
for i, txt in enumerate([
    "Read active HSD ID  (activeHsdId)",
    "POST /regression-check  �  Bridge Server � port 8776",
    "check_hsd_regression()  �  Call Intel HSD API  �  Fetch bug metadata",
], 1):
    pnode(s2, txt, MX, Y2[i], MW, BH, fill=_PROC_F, border=_PROC_B)

# 4: Decision � Versions
dnode(s2, "Versions complete?",
      MX, Y2[4], MW, BH, fill=_DEC_F, border=_DEC_B)

# 5: Auth check
pnode(s2, "/qb-check-auth  �  Verify QB session",
      MX, Y2[5], MW, BH, fill=_PROC_F, border=_PROC_B)

# 6: Decision � Logged in
dnode(s2, "QB session active?",
      MX, Y2[6], MW, BH, fill=_DEC_F, border=_DEC_B)

# 7-9: QB search steps
for i, txt in enumerate([
    "POST /qb-builds  �  qb_search_builds()",
    "Query QuickBuild API  �  Match fail / pass / fix builds",
    "Resolve CI build reference\nGFX: ci-master-NNNNN  |  GOP: gop-ci-main-NNN",
], 7):
    pnode(s2, txt, MX, Y2[i], MW, BH, fill=_PROC_F, border=_PROC_B)

# 10: End (filled green, white text)
pnode(s2, "Render Regression Table  in regression-area",
      MX, Y2[10], MW, BH, fill=_END_F, border=_END_F, bold=True, fc=WHITE)

# -- Vertical connectors (thin lines) ------------------------------------------
for i in range(10):
    rect(s2, MX+MW/2-0.01, Y2[i]+BH, 0.02, Y2[i+1]-Y2[i]-BH, fill=_LINE)

# -- Left bypass: "No" ? Version Input Overlay  (at Y2[4]) ---------------------
pnode(s2, "Show Version\nInput Overlay",
      BX, Y2[4], BW, BH, fill=_BYP_F, border=_BYP_B, fc=_TXT_D)
# horizontal connector decision ? branch
rect(s2, BX+BW, Y2[4]+BH/2-0.01, MX-(BX+BW), 0.02, fill=_LINE)
label(s2, "No", BX+BW+0.04, Y2[4]+BH/2-0.16, 0.28, 0.14,
      size=8, bold=True, color=_BYP_B)
# "Yes?" on main path
label(s2, "Yes?", MX-0.42, Y2[4]+BH+0.01, 0.40, 0.13,
      size=8, bold=True, color=_END_F, align=PP_ALIGN.RIGHT)
# merge connector
rect(s2, BX+BW, Y2[5], MX-(BX+BW), 0.02, fill=_LINE)

# -- Right bypass: "Not Logged In" ? QB Login Overlay  (at Y2[6]) --------------
pnode(s2, "Show QB Login Overlay\nEnter IDSID / Password",
      RX, Y2[6], 13.1-RX, BH, fill=_AUTH_F, border=_AUTH_B, fc=_TXT_D)
# horizontal connector decision ? branch
rect(s2, MX+MW, Y2[6]+BH/2-0.01, RX-(MX+MW), 0.02, fill=_LINE)
label(s2, "Not logged in", MX+MW+0.04, Y2[6]+BH/2-0.16, 0.88, 0.14,
      size=8, bold=True, color=_AUTH_B)
# "Logged in?" on main path
label(s2, "Logged in?", MX-0.72, Y2[6]+BH+0.01, 0.70, 0.13,
      size=8, bold=True, color=_END_F, align=PP_ALIGN.RIGHT)
# merge connector
rect(s2, MX+MW, Y2[7], RX-(MX+MW), 0.02, fill=_LINE)

# -- Show Commits  (horizontal, right of node 10) -------------------------------
rect(s2, MX+MW, Y2[10]+BH/2-0.01, RX-(MX+MW), 0.02, fill=_LINE)
dnode(s2, "Show\nCommits?",
      RX, Y2[10], 2.1, BH, fill=_DEC_F, border=_DEC_B, size=8)
rect(s2, RX+2.1, Y2[10]+BH/2-0.01, 0.22, 0.02, fill=_LINE)
label(s2, "Yes", RX+2.13, Y2[10]+BH/2-0.16, 0.22, 0.14,
      size=7.5, bold=True, color=_CMT_B)
pnode(s2, "POST /qb-commits  �  Trace revenue-pr\n? Commit table  +  View Changes on QB ?",
      RX+2.32, Y2[10], 13.1-(RX+2.32), BH,
      fill=_CMT_F, border=_CMT_B, size=8, fc=_TXT_D)


# ==============================================================================
# SLIDE 3  --  Phase 1 & Phase 2  (Professional English)
# ==============================================================================

s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, 13.33, 7.5, fill=WHITE)
header(s3, "Value Delivered", "Phase 1 & Phase 2", accent=INTEL_BLUE)
footer(s3, "GFX / GOP Driver Regression Checker � Value Proposition")

CARD_Y, CARD_H = 1.0, 5.85
_RED = RGBColor(0xDC, 0x26, 0x26)

# -- Phase 1 card (left) --------------------------------------------------------
rect(s3, 0.25, CARD_Y, 6.25, CARD_H,
     fill=RGBColor(0xEB, 0xEF, 0xFF), border=_PROC_B, bpt=1.5)
rect(s3, 0.25, CARD_Y, 6.25, 0.5, fill=_PROC_B)
label(s3, "Phase 1", 0.4, CARD_Y+0.03, 1.8, 0.44,
      size=22, bold=True, color=WHITE)
label(s3, "Automated CI Build Resolution",
      2.35, CARD_Y+0.12, 4.1, 0.32,
      size=11, color=RGBColor(0xC7, 0xD2, 0xFE))

# Problem
rect(s3, 0.35, 1.58, 6.05, 0.28, fill=_RED)
label(s3, "Problem  (Before)", 0.45, 1.60, 5.9, 0.24, size=9.5, bold=True, color=WHITE)
label(s3,
      "Production driver builds carry no metadata identifying the originating CI master branch.\n"
      "Engineers must manually search Driver History or QuickBuild to trace build lineage\n"
      "� time-consuming and error-prone.",
      0.45, 1.90, 5.95, 0.72, size=9.5, color=_TXT_D)

# Solution
rect(s3, 0.35, 2.70, 6.05, 0.28, fill=_END_F)
label(s3, "Solution  (After)", 0.45, 2.72, 5.9, 0.24, size=9.5, bold=True, color=WHITE)
label(s3,
      "Automatically resolves the corresponding ci-master build for any given\n"
      "Pass / Fail GFX / GOP driver version in under one minute.",
      0.45, 3.03, 5.95, 0.50, size=9.5, color=_TXT_D)

# Coverage
rect(s3, 0.35, 3.61, 6.05, 0.28, fill=RGBColor(0x1E, 0x29, 0x3B))
label(s3, "Current Coverage", 0.45, 3.63, 5.9, 0.24, size=9.5, bold=True, color=WHITE)

platforms = [
    ("GFX Driver",     "prod-hini-releases  ?  ci-master-NNNNN",                        _PROC_B),
    ("GOP  PTL",       "prod-gop-prod-PTL-N  /  prod-gop-Xe3-N  ?  gop-ci-main-NNN",   _DEC_B),
    ("GOP  WCL / NVL", "prod-gop-Xe3p-N  ?  ci-main-NNN",                              _DEC_B),
]
for i, (name, detail, col) in enumerate(platforms):
    py = 3.97 + i * 0.58
    rect(s3, 0.38, py, 0.08, 0.42, fill=col)
    rect(s3, 0.46, py, 5.8, 0.42,
         fill=WHITE, border=RGBColor(0xD1, 0xD5, 0xDB), bpt=0.5)
    label(s3, name,   0.57, py+0.04, 2.0, 0.18, size=9, bold=True, color=col)
    label(s3, detail, 0.57, py+0.22, 5.6, 0.18, size=8.5, color=DARK_GRAY)

# -- Phase 2 card (right) -------------------------------------------------------
rect(s3, 6.83, CARD_Y, 6.25, CARD_H,
     fill=RGBColor(0xF0, 0xFD, 0xF4), border=_END_F, bpt=1.5)
rect(s3, 6.83, CARD_Y, 6.25, 0.5, fill=_END_F)
label(s3, "Phase 2", 6.98, CARD_Y+0.03, 1.8, 0.44,
      size=22, bold=True, color=WHITE)
label(s3, "AI-Assisted Regression Isolation",
      8.9, CARD_Y+0.12, 4.1, 0.32,
      size=11, color=RGBColor(0xBB, 0xF7, 0xD0))

# Problem
rect(s3, 6.93, 1.58, 6.05, 0.28, fill=_RED)
label(s3, "Problem  (Before)", 7.03, 1.60, 5.9, 0.24, size=9.5, bold=True, color=WHITE)
label(s3,
      "Traditional bisection requires multiple manual validation cycles to\n"
      "identify the regression-inducing commit\n"
      "� labor-intensive, with no analytical guidance to prioritize candidates.",
      7.03, 1.90, 5.95, 0.72, size=9.5, color=_TXT_D)

# Solution
rect(s3, 6.93, 2.70, 6.05, 0.28, fill=_END_F)
label(s3, "Solution  (After)", 7.03, 2.72, 5.9, 0.24, size=9.5, bold=True, color=WHITE)
label(s3,
      "Integrates Sighting Assistant Tool analysis with CI build commit diffs to\n"
      "identify probable regression-inducing changes and generate a ranked list\n"
      "of build candidates � significantly reducing bisection overhead.",
      7.03, 3.03, 5.95, 0.72, size=9.5, color=_TXT_D)

# Steps
rect(s3, 6.93, 3.83, 6.05, 0.28, fill=RGBColor(0x1E, 0x29, 0x3B))
label(s3, "Execution Steps", 7.03, 3.85, 5.9, 0.24, size=9.5, bold=True, color=WHITE)

p2_steps = [
    ("?", "Retrieve commit delta between Pass / Fail ci-master builds",     _PROC_B),
    ("?", "Sighting Assistant Tool correlates HSD context with commit data", _END_F),
    ("?", "AI model ranks probable regression-inducing commits",             _END_F),
    ("?", "Output: prioritized build candidate list for targeted validation", _DEC_B),
]
for i, (num, step, col) in enumerate(p2_steps):
    sy = 4.19 + i * 0.52
    rect(s3, 6.95, sy, 0.34, 0.38, fill=col)
    label(s3, num, 6.95, sy, 0.34, 0.38,
          size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s3, 7.30, sy, 5.50, 0.38,
         fill=WHITE, border=RGBColor(0xD1, 0xD5, 0xDB), bpt=0.5)
    label(s3, step, 7.40, sy+0.07, 5.30, 0.26, size=9.5, color=_TXT_D)


# __ Save ______________________________________________________________________
out_path = r"C:\Users\brianch\Work\Chat_Mode_Assistant_tool\GFX_GOP_Regression_Checker_Overview.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
